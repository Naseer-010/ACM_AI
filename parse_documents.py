import io
import psycopg2
from datetime import datetime

from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from google.oauth2.credentials import Credentials
from google_auth import get_credentials

import pdfplumber
import docx
from pptx import Presentation

# =========================
# CONFIG
# =========================

DB_CONFIG = {
    "dbname": "studybuddy",
    "user": "nquzet",
    "host": "/var/run/postgresql",
}

SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]
TOKEN_PATH = "token.json"

# =========================
# DB CONNECTION
# =========================

conn = psycopg2.connect(**DB_CONFIG)
cursor = conn.cursor()

# =========================
# GOOGLE DRIVE CLIENT
# =========================

creds = get_credentials(SCOPES)
drive_service = build("drive", "v3", credentials=creds)

# =========================
# FILE DOWNLOAD
# =========================

def download_drive_file(file_id):
    request = drive_service.files().get_media(fileId=file_id)
    fh = io.BytesIO()
    downloader = MediaIoBaseDownload(fh, request)

    done = False
    while not done:
        _, done = downloader.next_chunk()

    fh.seek(0)
    return fh

# =========================
# PARSERS
# =========================

def parse_pdf(file_stream):
    text = []
    with pdfplumber.open(file_stream) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text.append(page_text)
    return "\n".join(text)


def parse_docx(file_stream):
    document = docx.Document(file_stream)
    return "\n".join(p.text for p in document.paragraphs if p.text.strip())


def parse_ppt(file_stream):
    prs = Presentation(file_stream)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text.strip():
                    text.append(shape.text)
    return "\n".join(text)

# =========================
# MAIN PIPELINE
# =========================

cursor.execute("""
    SELECT id, drive_file_id, file_type
    FROM documents
    WHERE parsed = FALSE
""")

documents = cursor.fetchall()
print(f"Found {len(documents)} unparsed documents")

for doc_id, drive_file_id, file_type in documents:
    print(f"Parsing document {doc_id} ({file_type})")

    try:
        file_stream = download_drive_file(drive_file_id)

        if file_type == "pdf":
            extracted_text = parse_pdf(file_stream)
        elif file_type == "docx":
            extracted_text = parse_docx(file_stream)
        elif file_type == "ppt":
            extracted_text = parse_ppt(file_stream)
        else:
            print(f"Unsupported file type: {file_type}")
            continue

        if not extracted_text.strip():
            print("No text extracted, skipping")
            continue

        cursor.execute("""
            UPDATE documents
            SET raw_text = %s,
                parsed = TRUE,
                created_at = created_at
            WHERE id = %s
        """, (
            extracted_text,
            doc_id
        ))

        conn.commit()
        print("✔ Parsed successfully")

    except Exception as e:
        conn.rollback()
        print(f"❌ Failed to parse document {doc_id}: {e}")

# =========================
# CLEANUP
# =========================

cursor.close()
conn.close()
print("Document parsing completed ✅")
